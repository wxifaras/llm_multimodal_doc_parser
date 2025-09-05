import os
import io
import base64
import json
from pathlib import Path
from pdf2image import convert_from_path
from openai import AzureOpenAI
import fitz
from PIL import Image
from pydantic import BaseModel
from PyPDF2 import PdfReader
from aoai import inference_structured_output_aoai, inference_aoai
from pptx import Presentation

try:
    import win32com.client  # For slide image export
except ImportError:
    win32com = None

aoai_deployment = os.environ.get("AOAI_DEPLOYMENT_NAME")
aoai_key = os.environ.get("AOAI_API_KEY")
aoai_endpoint = os.environ.get("AOAI_ENDPOINT")
aoai_api_version = os.environ.get("AOAI_API_VERSION")

pdf_path = os.environ.get("PDF_PATH")
ppt_path = os.environ.get("PPT_PATH")

document_structure_analyzer_prompt = """You are responsible for analyzing a document to determine its topic and top-level sections.

Output the following fields:
Summary: A brief description of the document type and its purpose. Make sure to include who/what the main subjects are.
Top-level Sections: What are the top-level sections of the document? 


###Examples###
User: (a large document of a home inspection report)
Assistant: 
Summary: Home Inspection Report for 337 Goldman Drive, Phoenixville, PA 19460, the home of Dan Giannone.
Top-level Sections: General Information
Introduction and Structural Overview
Purpose and Scope
Roof System
Exteriors
Electrical System
Heating and Cooling
Fireplaces and Chimneys
Plumbing System
Interiors
Basement and Crawlspace
Attic Area and Roof Framing
Appliances
NACHI Standards of Practice

"""

image_prompt = """You will be given an image that is one or more pages of a document, along with some analysis of the overall document. 

###Output Structure###
{
text: The verbatim text of the page in markdown format. 
images: A 1 sentence description of any images on the page and how they relate to the text.
image_insights: All insights or information that can be gleaned from the images on the page and the relationship to the text. 
}

###Guidance###

1. Output 'na' for images or image_insights if there are no images or diagrams. Do not consider tables images as you will be capturing them via text.
2. When outputting markdown, keep in mind that you are only looking at one page of a much larger document, so only consider something a section header if you feel very confident it could be a section header for this type of document.
3. In the text, make sure to indicate where the images are located on the page with an image tag like <image> </image>. You must describe the image and put that description within the image tag <image> </image>.
4. Use the surrounding text to provide context to the image & extract further insights from it. For example, if the text describes a picture of a house with "ADDRESS" listed below it, you can assume the image of the house is that address. Be as descriptive as possible. Just explain, do not start with "the image is..."
5. Only use markdown H2 headers for the top-level sections mentioned in the document structure analysis. Everything else should be a H3 or lower or some other markdown element.

###Examples###

User: (an image of the following text & picture) 
<document analysis>
Summary: Home Inspection Report for 337 Goldman Drive, Phoenixville, PA 19460, the home of Dan Giannone.
top-level sections: General Information
Introduction and Structural Overview
Purpose and Scope
Roof System
Exteriors
Electrical System
Heating and Cooling
Fireplaces and Chimneys
Plumbing System
Interiors
Basement and Crawlspace
Attic Area and Roof Framing
Appliances
NACHI Standards of Practice

<Content>
LDS Inspections
A Division of Lennox Design Studios

2801 Soni Drive Trooper, PA 19403
Phone: 610-277-4953 Fax: 610-277-4954
WWW.LDSINSPECTIONS.COM

Home Inspection Report For

---

(Image of a house)

337 Goldman Drive
Phoenixville, PA 19460

---

Report Prepared For
Dan Giannone

Report Prepared By
Craig Lennox


Assistant:

text: 
# LDS Inspections
**A Division of Lennox Design Studios**

2801 Soni Drive Trooper, PA 19403  
Phone: 610-277-4953 Fax: 610-277-4954  
[WWW.LDSINSPECTIONS.COM](http://www.ldsinspections.com)

**Home Inspection Report For**

---

[Image]

**337 Goldman Drive  
Phoenixville, PA 19460**

---

*Report Prepared For*  
*John Doe**

*Report Prepared By*  
**John Doe**

image_insights: 337 Goldman Dr, a large two-story suburban house owned by Dan Giannone. The house has the following features:

White exterior with light blue or gray trim
Multiple peaked roof sections
Several windows, including some arched windows on the upper floor
Two-car garage with white doors
Paved driveway with two vehicles parked in it (appear to be dark-colored sedans or similar)
Well-maintained front lawn
Some landscaping, including a small tree or bush with reddish foliage near the front of the house
Part of a neighboring house visible on the left side
Clear blue sky visible

"""

############################################################
# PDF processing functions
############################################################
def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")

def pdf_to_base64_images(pdf_path, output_dir):
    """
    Convert each PDF page to a PNG (in-memory) via PyMuPDF and return list of base64 strings.
    """
    doc = fitz.open(pdf_path)
    base64_images = []

    os.makedirs(output_dir, exist_ok=True)

    for page_index in range(len(doc)):
        page = doc.load_page(page_index)
        pix = page.get_pixmap()  # default zoom; adjust if higher resolution needed
        img = Image.open(io.BytesIO(pix.tobytes()))
        temp_image_path = os.path.join(output_dir, f"page_{page_index + 1}.png")
        img.save(temp_image_path, format="PNG")
        base64_images.append(encode_image(temp_image_path))
        #os.remove(temp_image_path) # uncomment to delete temp images

    return base64_images

class OutputStructure(BaseModel):
    text: str
    image_insights: str

def analyze_document_structure(pdf_path):
    reader = PdfReader(pdf_path)
    full_text = ""
    for page in reader.pages:
        full_text += page.extract_text()

    messages = [
        {"role": "system", "content": document_structure_analyzer_prompt},
        {"role": "user", "content": full_text}
    ]

    document_structure_analysis = inference_aoai(messages, aoai_deployment)
    return document_structure_analysis.choices[0].message.content

def process_image(image, page_num, source_filename, document_structure):
    messages = [
        {"role": "system", "content": f"{image_prompt}"},
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": (
                        f"source_filename: {source_filename}\n"
                        f"page_number: {page_num}\n\n"
                        f"Document Structure Analysis:\n{document_structure}"
                    )
                },
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image}", "detail": "high"}}
            ]
        }
    ]

    raw_response = inference_structured_output_aoai(messages, aoai_deployment, OutputStructure)
    if not raw_response:
        print(f"Failed to process page {page_num}")
        return None
    
    response = OutputStructure(**raw_response.choices[0].message.parsed.dict())
    print(f"Processed page {page_num}")
    print(f"Image Insights: {response.image_insights}")
    return response

def create_consolidated_markdown(processed_pages):
    consolidated_output = ""
    for page_num, page_data in enumerate(processed_pages, start=1):
        consolidated_output += page_data.text + "\n\n"
        if page_data.image_insights != 'na':
            consolidated_output += f"Image Insights: {page_data.image_insights}\n\n"
        consolidated_output += f"<Page {page_num}>\n\n"
        consolidated_output += "---\n\n"
    return consolidated_output

###########################################################
# PowerPoint processing functions
###########################################################
def extract_ppt_slide_texts(ppt_file: str) -> list[str]:
    """
    Extract plain text per slide using python-pptx.
    Returns list where index == slide_number - 1.
    """
    prs = Presentation(ppt_file)
    slide_texts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                txt = shape.text.strip()
                if txt:
                    parts.append(txt)
        # notes (optional)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(f"[Notes]\n{notes.strip()}")
        slide_texts.append("\n\n".join(parts) if parts else f"[Slide {idx}] (no extractable text)")
    return slide_texts

def export_ppt_slides_to_images(ppt_file: str, output_dir: str, width: int = 1920, height: int = 1080) -> list[str]:
    """
    Export each slide to a PNG using PowerPoint COM automation.
    Returns list of absolute file paths to exported PNGs in slide order.
    """
    if win32com is None:
        raise RuntimeError("pywin32 not installed (win32com.client). Install pywin32 or use a PDF fallback.")
    os.makedirs(output_dir, exist_ok=True)
    app = win32com.client.Dispatch("PowerPoint.Application")
    #app.Visible = 0
    pres = app.Presentations.Open(ppt_file, WithWindow=False)
    try:
        # PowerPoint's Export creates files: Slide1.PNG, Slide2.PNG, ...
        pres.Export(output_dir, "PNG", width, height)
        images = []
        for i in range(1, pres.Slides.Count + 1):
            candidate = os.path.join(output_dir, f"Slide{i}.PNG")
            if os.path.isfile(candidate):
                images.append(candidate)
        return images
    finally:
        pres.Close()
        app.Quit()

def slides_images_to_base64(image_paths: list[str]) -> list[str]:
    enc = []
    for p in image_paths:
        with open(p, "rb") as f:
            enc.append(base64.b64encode(f.read()).decode("utf-8"))
    return enc

def analyze_ppt_structure(ppt_file: str) -> str:
    """
    Build a 'document text' from all slides for structure inference (similar to PDF).
    """
    slide_texts = extract_ppt_slide_texts(ppt_file)
    combined = "\n\n---\n\n".join(f"[Slide {i+1}]\n{t}" for i, t in enumerate(slide_texts))
    messages = [
        {"role": "system", "content": document_structure_analyzer_prompt},
        {"role": "user", "content": combined}
    ]
    document_structure_analysis = inference_aoai(messages, aoai_deployment)
    return document_structure_analysis.choices[0].message.content, slide_texts

def process_slide(base64_image: str, slide_num: int, source_filename: str, document_structure: str, slide_text: str):
    """
    Mirrors process_image but includes raw slide text so the model can align visual & extracted text.
    """
    user_text = (
        f"source_filename: {source_filename}\n"
        f"slide_number: {slide_num}\n\n"
        f"Slide Text:\n{slide_text}\n\n"
        f"Document Structure Analysis:\n{document_structure}"
    )
    messages = [
        {"role": "system", "content": image_prompt},
        {
            "role": "user",
            "content": [
                {"type": "text", "text": user_text},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}", "detail": "high"}}
            ]
        }
    ]
    
    raw_response = inference_structured_output_aoai(messages, aoai_deployment, OutputStructure)
    if not raw_response:
        print(f"Failed to process slide {slide_num}")
        return None
    response = OutputStructure(**raw_response.choices[0].message.parsed.dict())
    print(f"Processed slide {slide_num}")
    print(f"Image Insights: {response.image_insights}")
    return response

def create_consolidated_markdown_slides(processed_slides: list[OutputStructure]):
    consolidated_output = ""
    for slide_num, slide_data in enumerate(processed_slides, start=1):
        consolidated_output += slide_data.text + "\n\n"
        if slide_data.image_insights != 'na':
            consolidated_output += f"Image Insights: {slide_data.image_insights}\n\n"
        consolidated_output += f"<Slide {slide_num}>\n\n"
        consolidated_output += "---\n\n"
    return consolidated_output

def process_ppt():
    if not os.path.isfile(ppt_path):
        raise FileNotFoundError(f"PPT/PPTX not found: {ppt_path}")
    directory = os.path.dirname(ppt_path)
    stem = Path(ppt_path).stem
    print(f"Using local PPTX: {ppt_path}")

    document_structure, slide_texts = analyze_ppt_structure(ppt_path)
    print("Presentation Structure Analysis:")
    print(document_structure)

    images_dir = os.path.join(directory, "slides_images")
    # Export slides to images
    try:
        image_paths = export_ppt_slides_to_images(ppt_path, images_dir)
    except RuntimeError as e:
        print(f"Image export failed ({e}); consider installing pywin32 / PowerPoint or fallback to PDF conversion.")
        image_paths = []

    if not image_paths:
        print("No slide images generated; cannot perform visual analysis. (You can convert PPTX to PDF and reuse PDF flow.)")
        return None

    base64_images = slides_images_to_base64(image_paths)

    processed_slides = []
    for slide_num, (b64, slide_text) in enumerate(zip(base64_images, slide_texts), start=1):
        slide_result = process_slide(b64, slide_num, os.path.basename(ppt_path), document_structure, slide_text)
        if slide_result:
            processed_slides.append(slide_result)

    consolidated = create_consolidated_markdown_slides(processed_slides)
    output_file = os.path.join(directory, f"{stem}_slides_consolidated.md")
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(consolidated)
    print(f"Slides processing complete. Consolidated results saved to {output_file}")
    return output_file

def process_pdf():
    """
    Process a local PDF (no blob storage). Assumes file exists at input_path/filename.
    """
    if not os.path.isfile(pdf_path):
        raise FileNotFoundError(f"PDF not found: {pdf_path}")
    
    directory = os.path.dirname(pdf_path)
    file_name = os.path.basename(pdf_path)
    stem = Path(file_name).stem

    print(f"Using local PDF: {pdf_path}")

    document_structure = analyze_document_structure(pdf_path)
    print("Document Structure Analysis:")
    print(document_structure)

    output_image_folder = f"{directory}\images"
    os.makedirs(output_image_folder, exist_ok=True)

    base64_images = pdf_to_base64_images(pdf_path, output_image_folder)

    processed_pages = []
    for page_num, image in enumerate(base64_images, start=1):
        page_result = process_image(image, page_num, file_name, document_structure)
        if page_result:
            processed_pages.append(page_result)

    consolidated_markdown = create_consolidated_markdown(processed_pages)

    output_file = os.path.join(directory, f"{stem}_consolidated.md")
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(consolidated_markdown)

    print(f"Processing complete. Consolidated results saved to {output_file}")
    return output_file

def main():
    """
    Auto-detect file type via environment variables:
      PDF_PATH -> process PDF
      PPT_PATH -> process PPTX/PPT
    Priority: PDF if both set.
    """
    if pdf_path and os.path.isfile(pdf_path):
        print("Detected PDF_PATH; running PDF pipeline.")
        return process_pdf()
    elif ppt_path and os.path.isfile(ppt_path):
        print("Detected PPT_PATH; running PPT pipeline.")
        return process_ppt()
    else:
        raise ValueError("Set PDF_PATH or PPT_PATH environment variable to a valid file.")


if __name__ == "__main__":

    main()